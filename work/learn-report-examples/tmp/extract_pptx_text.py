from pathlib import Path

from pptx import Presentation


pptx_path = Path("C:/ui5/mm.test-2026-06-141520/work/learn-report-examples/tmp/example-report.pptx")
out_path = Path("C:/ui5/mm.test-2026-06-141520/work/learn-report-examples/tmp/pptx-text.txt")

prs = Presentation(str(pptx_path))
lines = []
for slide_index, slide in enumerate(prs.slides, 1):
    lines.append(f"--- SLIDE {slide_index} ---")
    for shape_index, shape in enumerate(slide.shapes, 1):
        text = getattr(shape, "text", "")
        if text and text.strip():
            clean = "\n".join(part.rstrip() for part in text.splitlines() if part.strip())
            lines.append(f"[shape {shape_index}]")
            lines.append(clean)

out_path.write_text("\n".join(lines), encoding="utf-8")
print(out_path)
